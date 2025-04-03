import io
import json
from typing import Generator
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from aiedu.utils.image import image_compress, image_to_base64_url


def pptx_content_generator(
    pptx_path: str,
) -> Generator[tuple[list[str], list[Image.Image]], None, None]:
    """
    读取PPTX文件，获取文本、图片、表格和注释内容
    
    参数:
        pptx_path (str): PPTX文件路径
    
    返回:
        Generator[tuple[list[str], list[Image.Image]], None, None]: 生成器，返回文本、图片、表格和注释内容
    """
    # 获取PPTX文件
    presentation = Presentation(pptx_path)
    # 获取幻灯片
    for slide in presentation.slides:
        texts, images, tables, note = [], [], [], None
        for shape in slide.shapes:
            # 获取文本
            if hasattr(shape, "text"):
                texts.append(shape.text)
            # 获取图片
            if hasattr(shape, "image"):
                image = Image.open(io.BytesIO(shape.image.blob))
                image = image_compress(image)
                images.append(image_to_base64_url(image))
            # 获取表格
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = []
                for row in shape.table.rows:
                    table.append([cell.text for cell in row.cells])
                tables.append(json.dumps(table))
        # 获取注释
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
        # 生成器返回文本、图片、表格和注释
        yield texts, images, tables, note
